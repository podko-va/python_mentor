from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Создание презентации
prs = Presentation()

# Заголовки и содержимое для каждого слайда (на русском)
slides_content = [
    {
        "title": "Урок 4 — Обработка и агрегирование данных в Pandas",
        "content": "Сегодня мы изучим, как отбирать, группировать, объединять и трансформировать данные в библиотеке Pandas."
    },
    {
        "title": "Выбор данных (Data Selection)",
        "content": (
            "Методы:\n"
            "- .loc[] — выбор по метке\n"
            "- .iloc[] — выбор по позиции\n"
            "- .at[] и .iat[] — доступ к конкретной ячейке\n\n"
            "Примеры:\n"
            "df.loc[0:2, ['Name', 'Age']]\n"
            "df.iloc[:2]\n"
            "df[(df['Age'] > 24) & (df['Score'] >=88)]"
        )
    },
    {
        "title": "Агрегирование данных (Data Aggregation)",
        "content": (
            "Метод groupby() позволяет группировать данные по столбцам.\n\n"
            "Примеры:\n"
            "df.groupby('Category').sum()\n"
            "df.groupby('Category')['Values'].mean()\n"
            "df.groupby('Category').agg({'Values': ['sum', 'mean', 'count']})"
        )
    },
    {
        "title": "Объединение и соединение данных (Merging & Joining)",
        "content": (
            "Методы:\n"
            "- merge(): объединение по ключам\n"
            "- join(): соединение по индексам\n\n"
            "Примеры:\n"
            "pd.merge(df1, df2, on='ID', how='inner')\n"
            "df1.join(df2, how='outer')"
        )
    },
    {
        "title": "Преобразование данных (Data Transformation)",
        "content": (
            "Примеры операций:\n"
            "df['NewCol'] = df['OldCol'] ** 2\n"
            "df['EvenOdd'] = df['Value'].map(lambda x: 'Even' if x%2==0 else 'Odd')\n"
            "df['Value'] = df['Value'].astype(int)"
        )
    },
    {
        "title": "Утилиты Pandas (Полезные методы)",
        "content": (
            "Изменение имён столбцов:\n"
            "df.rename(columns={'Score':'Тест'}, inplace=True)\n\n"
            "Сортировка и индекс:\n"
            "df.sort_values(by='Score', ascending=False)\n"
            "df.set_index('Name')\n"
            "df.reset_index(drop=True)"
        )
    },
    {
        "title": "Контрольные вопросы",
        "content": (
            "- В чём разница между .loc и .iloc?\n"
            "- Как сгруппировать данные и вычислить среднее?\n"
            "- Чем отличается merge от join?\n"
            "- Как изменить или удалить столбец?\n"
            "- Какие методы сортировки и изменения индекса ты знаешь?"
        )
    }
]

# Добавление слайдов в презентацию
for slide in slides_content:
    slide_layout = prs.slide_layouts[1]
    s = prs.slides.add_slide(slide_layout)
    title = s.shapes.title
    content = s.placeholders[1]
    title.text = slide["title"]
    content.text = slide["content"]

# Сохранение презентации
pptx_path = "/mnt/data/Urok_04_Pandas_Data_Wrangling_Aggregation.pptx"
prs.save(pptx_path)

pptx_path
